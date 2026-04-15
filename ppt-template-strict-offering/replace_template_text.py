import argparse
import json
from pathlib import Path

from pptx import Presentation


def replace_shape_paragraphs(shape, lines):
    paragraphs = shape.text_frame.paragraphs
    for paragraph_index, paragraph in enumerate(paragraphs):
        text = lines[paragraph_index] if paragraph_index < len(lines) else ""
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = text


def main():
    parser = argparse.ArgumentParser(
        description="Copy a PPTX template and replace text in existing shape indexes."
    )
    parser.add_argument("--template", required=True, help="Reference PPTX path")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    parser.add_argument(
        "--replacements",
        required=True,
        help="JSON file mapping 1-based shape indexes to paragraph arrays",
    )
    parser.add_argument(
        "--slide",
        type=int,
        default=1,
        help="1-based slide index to edit (default: 1)",
    )
    args = parser.parse_args()

    template_path = Path(args.template)
    output_path = Path(args.output)
    replacements_path = Path(args.replacements)

    replacements = json.loads(replacements_path.read_text(encoding="utf-8"))

    presentation = Presentation(str(template_path))
    slide = presentation.slides[args.slide - 1]

    for shape_index_text, lines in replacements.items():
        shape_index = int(shape_index_text)
        shape = slide.shapes[shape_index - 1]

        if not hasattr(shape, "text_frame"):
            raise ValueError(f"Shape {shape_index} does not have a text frame")

        replace_shape_paragraphs(shape, lines)

    presentation.save(str(output_path))
    print(output_path)


if __name__ == "__main__":
    main()
